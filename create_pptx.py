"""
IC Candle Co. Business Plan 2026 - PowerPoint Generator
Dark theme: #0a0a0f background, white/light gray text, cyan #00e5ff and purple #a855f7 accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Constants ──────────────────────────────────────────────────────────
BG = RGBColor(10, 10, 15)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(232, 232, 236)
CYAN = RGBColor(0, 229, 255)
PURPLE = RGBColor(168, 85, 247)
DARK_CELL = RGBColor(18, 18, 28)
HEADER_CELL = RGBColor(28, 28, 42)
BORDER_COLOR = RGBColor(60, 60, 80)
DIM_GRAY = RGBColor(160, 160, 170)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ── Helper Functions ───────────────────────────────────────────────────

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, color=LIGHT_GRAY, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
                  font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_rich_paragraph(tf, runs, font_size=14, alignment=PP_ALIGN.LEFT,
                       space_before=Pt(4), space_after=Pt(2)):
    """Add a paragraph with multiple runs of different formatting.
    runs: list of (text, color, bold) tuples."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    for i, (text, color, bold) in enumerate(runs):
        if i == 0:
            r = p.runs[0] if p.runs else p.add_run()
        else:
            r = p.add_run()
        r.text = text
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return p


def add_slide_title(slide, title_text, top=Inches(0.4), left=Inches(0.8)):
    tf = add_text_box(slide, left, top, Inches(11.5), Inches(0.7),
                      title_text, font_size=32, color=WHITE, bold=True)
    return tf


def add_accent_line(slide, top, left=Inches(0.8), width=Inches(1.8)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    return shape


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            # Background
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = HEADER_CELL
            else:
                cell_fill.fore_color.rgb = DARK_CELL
            # Default text formatting
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def set_cell(table, row, col, text, color=LIGHT_GRAY, bold=False, size=11, align=PP_ALIGN.LEFT):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align


def set_col_widths(table, widths):
    for i, w in enumerate(widths):
        table.columns[i].width = w


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# Company name
add_text_box(slide, Inches(0), Inches(1.4), SLIDE_W, Inches(1.2),
             "IC CANDLE CO.", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Accent line centered
line_w = Inches(3)
line_left = Emu(int((SLIDE_W - line_w) / 2))
add_accent_line(slide, Inches(2.65), left=line_left, width=line_w)

# Subtitle
add_text_box(slide, Inches(0), Inches(2.85), SLIDE_W, Inches(0.6),
             "Business Plan 2026", font_size=28, color=CYAN, bold=False,
             alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(0), Inches(3.55), SLIDE_W, Inches(0.5),
             "Custom candles powered by AI.", font_size=18, color=DIM_GRAY,
             bold=False, alignment=PP_ALIGN.CENTER)

# Stats bar
add_text_box(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.5),
             "4 Members   |   $40/mo Operating Cost   |   $0 Marketing Agency",
             font_size=14, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 - EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Executive Summary")
add_accent_line(slide, Inches(1.05))

# Mission
tf = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), "", font_size=13)
p = tf.paragraphs[0]
p.text = "Mission"
p.font.size = Pt(18)
p.font.color.rgb = CYAN
p.font.bold = True

add_paragraph(tf, "To brighten homes with clean-burning, thoughtfully curated scents that provide a moment of peace \u2014 delivered through a modern, AI-driven e-commerce experience with zero waste and zero inventory risk.",
              font_size=13, color=LIGHT_GRAY, space_before=Pt(8))

add_paragraph(tf, "", font_size=8)

add_paragraph(tf, "Business Model", font_size=18, color=CYAN, bold=True, space_before=Pt(12))
add_paragraph(tf, "Customers order personalized candles from our Shopify storefront. Orders auto-route to our POD supplier who hand-pours and ships directly under our brand. AI handles all marketing \u2014 replacing $11K-19K/mo in agency costs.",
              font_size=13, color=LIGHT_GRAY, space_before=Pt(8))

add_paragraph(tf, "", font_size=8)

add_paragraph(tf, "Key Stats", font_size=18, color=CYAN, bold=True, space_before=Pt(12))
add_paragraph(tf, "\u2022  4-person team", font_size=13, color=LIGHT_GRAY, space_before=Pt(8))
add_paragraph(tf, "\u2022  ~$830 launch cost (~$208/person)", font_size=13)
add_paragraph(tf, "\u2022  Break-even at ~30 candles", font_size=13)
add_paragraph(tf, "\u2022  43-55% margins", font_size=13)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 - THE TEAM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "The Team \u2014 Four Roles. Zero Overlap.")
add_accent_line(slide, Inches(1.05))

team = [
    ("Ricardo", "Creative Director / Tech Lead",
     "Website design, AI prompt engineering, SEO content, brand identity, product mockups. Produces all creative assets with Claude."),
    ("Robert", "COO / Operations & Supply Chain",
     "Supplier sourcing, quality control, fulfillment, shipping, returns, customer service. Former REI head manager."),
    ("Jose", "Store Manager & Paid Marketing",
     "Shopify setup, product listings, Meta/TikTok ads, email platform (Klaviyo), analytics. Shopify experience."),
    ("Mark", "CFO & Business Strategy",
     "LLC formation, bookkeeping, P&L reporting, ad budget allocation, tax planning. Finance professional."),
]

y_pos = Inches(1.5)
for name, role, desc in team:
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), y_pos, Inches(11.5), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    tf = add_text_box(slide, Inches(1.1), y_pos + Inches(0.15), Inches(11), Inches(1.0), "", font_size=13)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = name
    r1.font.size = Pt(16)
    r1.font.color.rgb = CYAN
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = "  \u2014  " + role
    r2.font.size = Pt(14)
    r2.font.color.rgb = PURPLE
    r2.font.bold = True

    add_paragraph(tf, desc, font_size=12, color=LIGHT_GRAY, space_before=Pt(6))

    y_pos += Inches(1.4)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 - HOW THE MODEL WORKS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "How the Model Works")
add_accent_line(slide, Inches(1.05))

steps = [
    ("1", "Customer Orders", "Buyer places order on Shopify/Etsy"),
    ("2", "Auto-Routed", "Order routes to POD supplier via app"),
    ("3", "Produced", "Supplier hand-pours with your label"),
    ("4", "Shipped", "Ships to customer under your brand"),
    ("5", "You Profit", "Zero inventory, zero touchpoints"),
]

x_start = Inches(0.5)
step_w = Inches(2.3)
gap = Inches(0.2)
y_top = Inches(2.0)

for i, (num, title, desc) in enumerate(steps):
    x = x_start + i * (step_w + gap)

    # Step box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, y_top, step_w, Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_CELL
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(0.5)

    # Step number
    add_text_box(slide, x, y_top + Inches(0.2), step_w, Inches(0.6),
                 num, font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    # Step title
    add_text_box(slide, x + Inches(0.15), y_top + Inches(0.9), step_w - Inches(0.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Step description
    add_text_box(slide, x + Inches(0.15), y_top + Inches(1.5), step_w - Inches(0.3), Inches(1.0),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        arrow_x = x + step_w + Inches(0.02)
        add_text_box(slide, arrow_x, y_top + Inches(1.0), Inches(0.2), Inches(0.4),
                     "\u25B6", font_size=14, color=CYAN, alignment=PP_ALIGN.CENTER)

# Key insight
tf = add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8), "", font_size=14)
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "Key Insight: "
r1.font.size = Pt(15)
r1.font.color.rgb = CYAN
r1.font.bold = True
r2 = p.add_run()
r2.text = "You only pay the supplier AFTER the customer pays you."
r2.font.size = Pt(15)
r2.font.color.rgb = WHITE
r2.font.bold = False
p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 - OUR UNFAIR ADVANTAGE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Our Unfair Advantage")
add_accent_line(slide, Inches(1.05))

roles_data = [
    ("Web Designer", "$2,000-4,000/mo", "$0"),
    ("SEO Specialist", "$1,500-3,000/mo", "$0"),
    ("Copywriter", "$2,000-3,000/mo", "$0"),
    ("Email Marketer", "$1,000-2,000/mo", "$0"),
    ("Social Media Manager", "$1,500-2,000/mo", "$0"),
    ("Brand Strategist", "$2,000-3,000/mo", "$0"),
    ("Market Researcher", "$1,500-2,000/mo", "$0"),
]

rows = len(roles_data) + 2  # header + data + total
table = add_table(slide, rows, 3, Inches(1.5), Inches(1.5), Inches(10), Inches(5.0))

set_col_widths(table, [Inches(4), Inches(3.5), Inches(2.5)])

set_cell(table, 0, 0, "Agency Role", CYAN, True, 13, PP_ALIGN.LEFT)
set_cell(table, 0, 1, "Market Rate", CYAN, True, 13, PP_ALIGN.CENTER)
set_cell(table, 0, 2, "Our Cost", CYAN, True, 13, PP_ALIGN.CENTER)

for i, (role, rate, cost) in enumerate(roles_data, 1):
    set_cell(table, i, 0, role, LIGHT_GRAY, False, 12)
    set_cell(table, i, 1, rate, LIGHT_GRAY, False, 12, PP_ALIGN.CENTER)
    set_cell(table, i, 2, cost, RGBColor(0, 255, 100), True, 13, PP_ALIGN.CENTER)

# Total row
total_row = len(roles_data) + 1
set_cell(table, total_row, 0, "TOTAL", WHITE, True, 13)
set_cell(table, total_row, 1, "$11,500-19,000/mo", PURPLE, True, 13, PP_ALIGN.CENTER)
set_cell(table, total_row, 2, "$0", RGBColor(0, 255, 100), True, 14, PP_ALIGN.CENTER)

# Highlight total row
for col_idx in range(3):
    cell = table.cell(total_row, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(30, 30, 50)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 - MARKET ANALYSIS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Market Analysis")
add_accent_line(slide, Inches(1.05))

sections = [
    ("Industry Trends (2026)", [
        "Global scented candle market growing",
        "Shift to sustainability and wellness",
    ]),
    ("Target Market", [
        "Age 25-40, homeowners/apartment dwellers",
        "Eco-conscious, self-care focused",
    ]),
    ("Competitive Advantage", [
        "Unique artisanal scents, non-toxic wax",
        "Lean POD model eliminates inventory risk",
    ]),
]

y_pos = Inches(1.5)
for section_title, bullets in sections:
    # Section card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), y_pos, Inches(11.5), Inches(1.6))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    tf = add_text_box(slide, Inches(1.1), y_pos + Inches(0.15), Inches(11), Inches(1.4), "", font_size=13)
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(17)
    p.font.color.rgb = CYAN
    p.font.bold = True

    for bullet in bullets:
        add_paragraph(tf, "\u2022  " + bullet, font_size=13, color=LIGHT_GRAY, space_before=Pt(6))

    y_pos += Inches(1.85)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 - NICHE & POSITIONING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Niche & Positioning")
add_accent_line(slide, Inches(1.05))

# Principle
tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), "", font_size=14)
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "The scent is not the niche \u2014 the "
r1.font.size = Pt(15)
r1.font.color.rgb = LIGHT_GRAY
r2 = p.add_run()
r2.text = "STORY"
r2.font.size = Pt(15)
r2.font.color.rgb = CYAN
r2.font.bold = True
r3 = p.add_run()
r3.text = " around it is."
r3.font.size = Pt(15)
r3.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

niche_data = [
    ("Vanilla / Warm Spice", "Wedding Favors, Bridal Gifts", "Brides, Wedding Planners"),
    ("Eucalyptus / Lavender", "Wellness, Yoga, Self-Care", "Wellness Community"),
    ("Sea Breeze / Ocean", "Coastal Decor, Beach Homes", "Coastal Lifestyle Buyers"),
    ("Cedar / Woodsy", "Masculine Gifting, Outdoors", "Men's Gift Market"),
    ("Citrus / Grapefruit", "Energizing, Morning Routines", "Productivity & Wellness"),
    ("Pumpkin / Balsam", "Seasonal / Holiday", "Holiday Gift Buyers"),
]

table = add_table(slide, len(niche_data) + 1, 3,
                  Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.8))

set_col_widths(table, [Inches(3.8), Inches(4.0), Inches(3.7)])

set_cell(table, 0, 0, "Supplier Scent", CYAN, True, 12)
set_cell(table, 0, 1, "Niche Angle", CYAN, True, 12)
set_cell(table, 0, 2, "Target Buyer", CYAN, True, 12)

for i, (scent, angle, buyer) in enumerate(niche_data, 1):
    set_cell(table, i, 0, scent, LIGHT_GRAY, False, 12)
    set_cell(table, i, 1, angle, LIGHT_GRAY, False, 12)
    set_cell(table, i, 2, buyer, LIGHT_GRAY, False, 12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 - SALES CHANNELS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Sales Channels \u2014 Sell Everywhere")
add_accent_line(slide, Inches(1.05))

channels = [
    ("Shopify", "Home base"),
    ("Etsy", "#1 for personalized gifts"),
    ("TikTok Shop", "Sell in viral videos"),
    ("Instagram Shop", "Tag in posts/reels"),
    ("Amazon", "Massive trust + gift traffic"),
    ("SEO Blog", "Free Google traffic"),
]

# 2 rows of 3 cards
for idx, (ch_name, ch_desc) in enumerate(channels):
    row = idx // 3
    col = idx % 3
    x = Inches(0.8) + col * Inches(4.0)
    y = Inches(1.6) + row * Inches(2.8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(3.6), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.2), Inches(0.6),
                 ch_name, font_size=22, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.2), Inches(0.8),
                 ch_desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 - MARGIN MODEL
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Margin Model \u2014 Mid-Range Pricing ($25-$45)")
add_accent_line(slide, Inches(1.05))

# Two columns
for col_idx, (label, data) in enumerate([
    ("Conservative", [
        ("Retail Price", "$28"),
        ("Supplier Cost", "~$13"),
        ("Fees (10%)", "~$2.80"),
        ("Margin", "~$12"),
        ("Margin %", "~43%"),
    ]),
    ("Strong", [
        ("Retail Price", "$42"),
        ("Supplier Cost", "~$15"),
        ("Fees (10%)", "~$4.20"),
        ("Margin", "~$23"),
        ("Margin %", "~55%"),
    ]),
]):
    x = Inches(0.8) + col_idx * Inches(6.0)

    # Column header
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(1.5), Inches(5.5), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = HEADER_CELL
    box.line.color.rgb = BORDER_COLOR

    accent = CYAN if col_idx == 0 else PURPLE
    add_text_box(slide, x, Inches(1.55), Inches(5.5), Inches(0.5),
                 label, font_size=18, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    table = add_table(slide, len(data), 2, x, Inches(2.2), Inches(5.5), Inches(3.0))
    set_col_widths(table, [Inches(3.0), Inches(2.5)])
    for i, (item, val) in enumerate(data):
        set_cell(table, i, 0, item, LIGHT_GRAY, False, 13)
        val_color = RGBColor(0, 255, 100) if "Margin" in item else WHITE
        set_cell(table, i, 1, val, val_color, True if "Margin" in item else False, 13, PP_ALIGN.RIGHT)

# Shipping strategies
tf = add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2), "", font_size=13)
p = tf.paragraphs[0]
p.text = "Shipping Strategies"
p.font.size = Pt(16)
p.font.color.rgb = CYAN
p.font.bold = True
add_paragraph(tf, "\u2022  Pass to customer   \u2022  Build into price   \u2022  Free threshold (2+ candles)",
              font_size=13, color=LIGHT_GRAY, space_before=Pt(8))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 - SUPPLIER SCORECARD (Operations)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Supplier Scorecard \u2014 Product, Fulfillment & Integration")
add_accent_line(slide, Inches(1.05))

# Headers
headers = ["", "Candle Bliss", "Candle Bldrs", "Printify", "Printed Mint", "Printful"]

scorecard_data = [
    # Product section
    ("PRODUCT", None),
    ("True POD", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("White Label", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("Scent Variety", ["\u2705 20+", "\u2705 8+", "\u2705 10+", "\u26a0 4", "\u26a0 9"]),
    ("Vessel Options", ["\u2705 Tins+3", "\u26a0 3 glass", "\u26a0 Glass", "\u2705 Multi", "\u274c 1 jar"]),
    ("Wax Type", ["Coco/Apricot", "Soy", "Soy", "Coco/Apricot", "Soy"]),
    ("Gift Packaging", ["\u2705", "\u2705 WL", "\u274c", "\u2705 Full", "\u274c"]),
    # Fulfillment section
    ("FULFILLMENT", None),
    ("Speed", ["\u2705 1-3d", "\u2705 1-2d", "\u26a0 Varies", "\u26a0 4-7d", "\u2705 Fast"]),
    ("US-Based", ["\u2705 AZ", "\u2705 NH", "\u2705 Most", "\u2705", "\u2705"]),
    ("No MOQ", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    # Integration section
    ("INTEGRATION", None),
    ("Shopify", ["\u2705 Native", "\u2705 Native", "\u2705", "\u2705", "\u2705"]),
    ("Etsy", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("Amazon", ["\u274c", "\u274c", "\u2705 Multi", "\u2705 OD", "\u2705 23+"]),
    ("Personalization", ["\u274c", "\u274c", "\u274c", "\u274c", "\u274c"]),
    ("Startup Cost", ["\u2705 Free", "\u2705 Free", "\u2705 Free", "\u2705 Free", "\u2705 Free"]),
]

# Count actual data rows (not section headers)
data_rows = [r for r in scorecard_data if r[1] is not None]
section_rows = [r for r in scorecard_data if r[1] is None]

total_rows = len(scorecard_data) + 1  # +1 for header
table = add_table(slide, total_rows, 6, Inches(0.3), Inches(1.25), Inches(12.6), Inches(5.8))

set_col_widths(table, [Inches(1.8), Inches(2.16), Inches(2.16), Inches(2.16), Inches(2.16), Inches(2.16)])

# Set header row
for i, h in enumerate(headers):
    set_cell(table, 0, i, h, CYAN, True, 10, PP_ALIGN.CENTER)

row_idx = 1
for label, values in scorecard_data:
    if values is None:
        # Section header row
        cell = table.cell(row_idx, 0)
        # Merge across all columns for section header
        for ci in range(6):
            c = table.cell(row_idx, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(25, 25, 40)
        set_cell(table, row_idx, 0, label, PURPLE, True, 10, PP_ALIGN.LEFT)
        for ci in range(1, 6):
            set_cell(table, row_idx, ci, "", LIGHT_GRAY, False, 10)
    else:
        set_cell(table, row_idx, 0, label, LIGHT_GRAY, False, 10)
        for ci, val in enumerate(values, 1):
            set_cell(table, row_idx, ci, val, LIGHT_GRAY, False, 9, PP_ALIGN.CENTER)
    row_idx += 1


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 - SUPPLIER SCORECARD (Sustainability)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Supplier Scorecard \u2014 Sustainability & Risk")
add_accent_line(slide, Inches(1.05))

headers2 = ["", "Candle Bliss", "Candle Bldrs", "Printify", "Printed Mint", "Printful"]

sust_data = [
    ("SUSTAINABILITY", None),
    ("Plant-Based Wax", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("No Paraffin", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("Vegan/Cruelty-Free", ["\u2705", "\u2705", "\u2705", "\u2705", "\u26a0 N/S"]),
    ("Phthalate-Free", ["\u2705 Prop65", "\u2705", "\u26a0 Varies", "\u2705", "\u26a0 N/S"]),
    ("Clean Wick", ["\u26a0 N/S", "\u2705 Cotton", "\u26a0 Varies", "\u26a0 N/S", "\u26a0 N/S"]),
    ("No Synthetic Dyes", ["\u2705", "\u26a0 N/S", "\u26a0 Varies", "\u26a0 N/S", "\u26a0 N/S"]),
    ("Recyclable Vessel", ["\u2705", "\u2705", "\u2705", "\u2705", "\u2705"]),
    ("USA-Sourced", ["\u2705 Wax+Fr", "\u2705 Wax+Fr", "\u26a0 Varies", "\u26a0 N/S", "\u26a0 N/S"]),
    ("RISK", None),
    ("Record Keeping", ["\u2705 Auto", "\u2705 Auto", "\u2705 Dashboard", "\u26a0 Manual", "\u2705 Full"]),
]

total_rows2 = len(sust_data) + 1
table2 = add_table(slide, total_rows2, 6, Inches(0.3), Inches(1.25), Inches(12.6), Inches(5.0))
set_col_widths(table2, [Inches(1.8), Inches(2.16), Inches(2.16), Inches(2.16), Inches(2.16), Inches(2.16)])

for i, h in enumerate(headers2):
    set_cell(table2, 0, i, h, CYAN, True, 10, PP_ALIGN.CENTER)

row_idx = 1
for label, values in sust_data:
    if values is None:
        for ci in range(6):
            c = table2.cell(row_idx, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(25, 25, 40)
        set_cell(table2, row_idx, 0, label, PURPLE, True, 10, PP_ALIGN.LEFT)
        for ci in range(1, 6):
            set_cell(table2, row_idx, ci, "", LIGHT_GRAY, False, 10)
    else:
        set_cell(table2, row_idx, 0, label, LIGHT_GRAY, False, 10)
        for ci, val in enumerate(values, 1):
            set_cell(table2, row_idx, ci, val, LIGHT_GRAY, False, 9, PP_ALIGN.CENTER)
    row_idx += 1

# Footnote
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             "N/S = Not stated publicly", font_size=11, color=DIM_GRAY,
             alignment=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 - RECOMMENDED SUPPLIERS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Recommended Suppliers")
add_accent_line(slide, Inches(1.05))

suppliers = [
    ("#1  Candle Bliss", "candlebliss.com",
     "Best scent library (20+), strong packaging, free apps, Phoenix AZ, 1-3 day turnaround.",
     "Best for: gifting, scent-forward.",
     "Gap: Coco/apricot only."),
    ("#2  Printed Mint", "printedmint.com",
     "Best gift packaging.",
     "Best for: gift niche, unboxing.",
     "Gap: 4 scents, 4-7 days, manual records."),
    ("#3  Candle Builders", "candlebuilders.com",
     "Fast (1-2 days), native integrations, 8+ scents.",
     "Best for: speed, soy positioning.",
     "Gap: Soy only, 3 glass types."),
]

y = Inches(1.5)
for rank_name, url, desc, best_for, gap in suppliers:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), y, Inches(11.5), Inches(1.45))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    tf = add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(11), Inches(1.3), "", font_size=12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = rank_name
    r1.font.size = Pt(16)
    r1.font.color.rgb = CYAN
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = "  (" + url + ")"
    r2.font.size = Pt(12)
    r2.font.color.rgb = DIM_GRAY

    add_paragraph(tf, desc, font_size=12, color=LIGHT_GRAY, space_before=Pt(4))
    add_paragraph(tf, best_for, font_size=12, color=RGBColor(0, 255, 100), space_before=Pt(2))
    add_paragraph(tf, gap, font_size=11, color=DIM_GRAY, space_before=Pt(2))

    y += Inches(1.65)

# Recommendation
rec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), y + Inches(0.15), Inches(11.5), Inches(0.7))
rec_box.fill.solid()
rec_box.fill.fore_color.rgb = RGBColor(25, 25, 40)
rec_box.line.color.rgb = CYAN
rec_box.line.width = Pt(1)

tf = add_text_box(slide, Inches(1.1), y + Inches(0.25), Inches(11), Inches(0.5), "", font_size=14)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Recommendation: "
r.font.size = Pt(14)
r.font.color.rgb = CYAN
r.font.bold = True
r2 = p.add_run()
r2.text = "Launch with Candle Bliss on Etsy. Add Printed Mint for gift SKUs later."
r2.font.size = Pt(14)
r2.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 - THE NUMBERS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "The Numbers")
add_accent_line(slide, Inches(1.05))

# Three columns
col_data = [
    ("Startup Costs", CYAN, [
        "FL LLC: $125",
        "Shopify: $39",
        "Domain: $15",
        "Samples: $150",
        "Ads: $500",
        "",
        "Total: ~$830",
        "(~$208/person)",
    ]),
    ("Unit Economics", PURPLE, [
        "Sale: $28-38",
        "POD: -$12-16",
        "Shipping: -$5-8",
        "Profit: $10-14",
        "",
        "Overhead: ~$140/mo",
        "Break-even: ~30 candles",
    ]),
    ("Profit Split", CYAN, [
        "60% Reinvest",
        "25% Reserve",
        "15% Member Distribution",
    ]),
]

for ci, (col_title, accent, items) in enumerate(col_data):
    x = Inches(0.6) + ci * Inches(4.2)
    w = Inches(3.8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.5), w, Inches(5.2))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    add_text_box(slide, x + Inches(0.15), Inches(1.65), w - Inches(0.3), Inches(0.5),
                 col_title, font_size=18, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + Inches(0.3), Inches(2.3), w - Inches(0.6), Inches(4.2),
                      "", font_size=13)
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Calibri"
            if "Total" in item or "Profit:" in item or "Break-even" in item:
                p.font.color.rgb = WHITE
                p.font.bold = True
            first = False
        else:
            color = WHITE if ("Total" in item or "Profit:" in item or "Break-even" in item) else LIGHT_GRAY
            bld = "Total" in item or "Profit:" in item or "Break-even" in item
            add_paragraph(tf, item, font_size=13, color=color, bold=bld, space_before=Pt(6))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 - FOUNDATION CHECKLIST (Weeks 1-2)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Foundation \u2014 Weeks 1-2")
add_accent_line(slide, Inches(1.05))

checklist_14 = [
    ("Ricardo", CYAN, [
        "Research competitor stores",
        "Select color palette",
        "Design logo",
        "Homepage mockup",
        "3 label templates",
        "Set up IG + TikTok",
        "Brand style guide",
    ]),
    ("Robert", PURPLE, [
        "Compare top 5 suppliers",
        "Apply Candle Bliss",
        "Apply Printed Mint",
        "Order samples",
        "Document terms",
        "Research insurance",
        "Supplier spreadsheet",
    ]),
    ("Jose", CYAN, [
        "Create Shopify ($39)",
        "Buy domain",
        "Install Candle Bliss + Klaviyo",
        "Shopify Payments + PayPal",
        "Research Etsy pricing",
        "GA + Pixel",
        "Placeholder collections",
    ]),
    ("Mark", PURPLE, [
        "File FL LLC ($125 Sunbiz)",
        "EIN (IRS.gov)",
        "Bank account (Mercury)",
        "Wave accounting",
        "Operating agreement",
        "Budget tracker",
        "FL sales tax research",
    ]),
]

for ci, (name, accent, tasks) in enumerate(checklist_14):
    x = Inches(0.3) + ci * Inches(3.2)
    w = Inches(3.0)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.4), w, Inches(5.6))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    add_text_box(slide, x + Inches(0.1), Inches(1.5), w - Inches(0.2), Inches(0.5),
                 name, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + Inches(0.15), Inches(2.1), w - Inches(0.3), Inches(4.8),
                      "", font_size=11)
    first = True
    for task in tasks:
        if first:
            p = tf.paragraphs[0]
            p.text = "\u25a1  " + task
            p.font.size = Pt(11)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Calibri"
            first = False
        else:
            add_paragraph(tf, "\u25a1  " + task, font_size=11, color=LIGHT_GRAY, space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 - BUILD CHECKLIST (Weeks 3-4)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Build \u2014 Weeks 3-4")
add_accent_line(slide, Inches(1.05))

checklist_15 = [
    ("Ricardo", CYAN, [
        "Full Shopify theme",
        "Product pages",
        "Customizer tool",
        "AI mockups",
        "SEO descriptions (8-10)",
        "Email popup",
        "Social templates",
    ]),
    ("Robert", PURPLE, [
        "Burn-test samples",
        "Select supplier",
        "Set up fulfillment app",
        "Shipping rates",
        "Returns policy",
        "CS templates (10)",
        "Secure insurance",
    ]),
    ("Jose", CYAN, [
        "Upload all products",
        "Shipping + tax config",
        "Open Etsy",
        "5 Etsy listings",
        "Klaviyo capture",
        "Welcome sequence (3 emails)",
        "Meta Business Manager",
    ]),
    ("Mark", PURPLE, [
        "Pricing model",
        "Unit economics",
        "Weekly P&L template",
        "Sales tax collection",
        "Ad budget rules",
        "3/6/12-month projections",
        "Profit-split tracking",
    ]),
]

for ci, (name, accent, tasks) in enumerate(checklist_15):
    x = Inches(0.3) + ci * Inches(3.2)
    w = Inches(3.0)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.4), w, Inches(5.6))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    add_text_box(slide, x + Inches(0.1), Inches(1.5), w - Inches(0.2), Inches(0.5),
                 name, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + Inches(0.15), Inches(2.1), w - Inches(0.3), Inches(4.8),
                      "", font_size=11)
    first = True
    for task in tasks:
        if first:
            p = tf.paragraphs[0]
            p.text = "\u25a1  " + task
            p.font.size = Pt(11)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Calibri"
            first = False
        else:
            add_paragraph(tf, "\u25a1  " + task, font_size=11, color=LIGHT_GRAY, space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 - LAUNCH, ADS & SCALE (Weeks 5-8+)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Launch, Ads & Scale \u2014 Weeks 5-8+")
add_accent_line(slide, Inches(1.05))

launch_sections = [
    ("Soft Launch (Wk 5-6)", CYAN, [
        "QA store", "Launch graphics", "Blog posts",
        "Daily social", "Test orders", "Verify pipeline",
        "GO LIVE", "Instagram Shopping", "TikTok Shop",
        "Collect reviews", "Record financials", "Verify tax",
    ]),
    ("Paid Ads (Wk 7-8)", PURPLE, [
        "5 ad creatives", "Retargeting", "More blog posts",
        "3 Meta campaigns ($20-30/day)", "Abandoned cart flow",
        "TikTok push", "Track ROAS", "Calculate CAC",
        "Monthly P&L",
    ]),
    ("Scale (Mo 3+)", CYAN, [
        "Double winning ads", "Google Shopping", "TikTok ads",
        "Seasonal collection", "Amazon", "Add Printed Mint",
        "Reinvest 60%", "Quarterly review",
    ]),
]

for ci, (section_title, accent, items) in enumerate(launch_sections):
    x = Inches(0.4) + ci * Inches(4.2)
    w = Inches(3.9)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.4), w, Inches(5.6))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    add_text_box(slide, x + Inches(0.1), Inches(1.5), w - Inches(0.2), Inches(0.5),
                 section_title, font_size=15, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + Inches(0.15), Inches(2.1), w - Inches(0.3), Inches(4.8),
                      "", font_size=11)
    first = True
    for item in items:
        is_highlight = item == "GO LIVE"
        color = WHITE if is_highlight else LIGHT_GRAY
        bld = is_highlight
        if first:
            p = tf.paragraphs[0]
            p.text = "\u25b8  " + item
            p.font.size = Pt(11)
            p.font.color.rgb = color
            p.font.bold = bld
            p.font.name = "Calibri"
            first = False
        else:
            add_paragraph(tf, "\u25b8  " + item, font_size=11, color=color, bold=bld, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 - RISK MANAGEMENT
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_title(slide, "Risk Management")
add_accent_line(slide, Inches(1.05))

risk_sections = [
    ("Chargeback Risk", [
        "Seller of record liability",
        "Document everything",
    ]),
    ("Returns", [
        "No supplier returns on POD",
        "Replace damaged",
        "Store credit for changes",
        "Final sale on personalized",
    ]),
    ("Insurance", [
        "General + Product Liability $300-600/yr",
        "Next / Thimble / Hiscox",
        "Get before first sale",
    ]),
    ("Record Keeping", [
        "Log every order",
        "Screenshot comms",
        "Use dashboards",
        "Manual for Printed Mint",
    ]),
]

# 2x2 grid
for idx, (section_title, bullets) in enumerate(risk_sections):
    row = idx // 2
    col = idx % 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(1.5) + row * Inches(2.8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(5.5), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CELL
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(0.5)

    tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.1), Inches(2.2),
                      "", font_size=13)
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(17)
    p.font.color.rgb = CYAN
    p.font.bold = True

    for bullet in bullets:
        add_paragraph(tf, "\u2022  " + bullet, font_size=12, color=LIGHT_GRAY, space_before=Pt(6))


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_path = "IC_Candle_Co_Business_Plan_2026.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
